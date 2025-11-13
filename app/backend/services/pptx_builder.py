from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Dict, List
from config import Config
import os

class PPTXBuilder:
    """Build PowerPoint presentations from slide data"""
    
    THEME_COLORS = {
        'blue': {
            'primary': RGBColor(30, 64, 175),
            'secondary': RGBColor(59, 130, 246),
            'accent': RGBColor(219, 234, 254)
        },
        'green': {
            'primary': RGBColor(22, 101, 52),
            'secondary': RGBColor(34, 197, 94),
            'accent': RGBColor(220, 252, 231)
        },
        'purple': {
            'primary': RGBColor(107, 33, 168),
            'secondary': RGBColor(168, 85, 247),
            'accent': RGBColor(243, 232, 255)
        },
        'orange': {
            'primary': RGBColor(194, 65, 12),
            'secondary': RGBColor(249, 115, 22),
            'accent': RGBColor(254, 215, 170)
        }
    }
    
    def __init__(self, theme: str = 'blue'):
        self.prs = Presentation()
        self.prs.slide_width = Inches(Config.SLIDE_WIDTH)
        self.prs.slide_height = Inches(Config.SLIDE_HEIGHT)
        self.theme = self.THEME_COLORS.get(theme, self.THEME_COLORS['blue'])
    
    def create_presentation(self, slides_data: List[Dict], output_path: str) -> str:
        """Create complete presentation from slide data"""
        
        for slide_data in slides_data:
            slide_type = slide_data.get('type')
            content = slide_data.get('content', {})
            
            if slide_type == 'title':
                self._add_title_slide(content)
            elif slide_type == 'executive':
                self._add_executive_slide(content)
            elif slide_type == 'financials':
                self._add_financials_slide(content)
            elif slide_type in ['assets', 'liabilities']:
                self._add_breakdown_slide(content, slide_type)
            elif slide_type == 'ratios':
                self._add_ratios_slide(content)
            elif slide_type in ['trends', 'conclusion']:
                self._add_insights_slide(content)
            elif slide_type == 'company':
                self._add_company_slide(content)
            elif slide_type == 'products_services':
                self._add_products_services_slide(content)
            elif slide_type == 'markets_locations':
                self._add_markets_locations_slide(content)
            elif slide_type == 'leadership':
                self._add_leadership_slide(content)
            elif slide_type == 'major_projects':
                self._add_major_projects_slide(content)
            elif slide_type == 'vision_mission':
                self._add_vision_mission_slide(content)
        
        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Save presentation
        self.prs.save(output_path)
        return output_path
    
    def _add_title_slide(self, content: Dict):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add background
        background = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(8), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Financial Presentation')
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(3.7),
            Inches(8), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = content.get('subtitle', '')
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Add company name
        company_box = slide.shapes.add_textbox(
            Inches(1), Inches(5),
            Inches(8), Inches(0.8)
        )
        company_frame = company_box.text_frame
        company_frame.text = content.get('company_name', '')
        company_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        company_frame.paragraphs[0].font.size = Pt(24)
        company_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    
    def _add_executive_slide(self, content: Dict):
        """Add executive summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
            )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Executive Summary')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Highlights
        highlights = content.get('highlights', [])
        y_position = 1.5
        
        for i, highlight in enumerate(highlights[:5]):
            highlight_box = slide.shapes.add_textbox(
                Inches(1), Inches(y_position),
                Inches(8), Inches(0.8)
            )
            tf = highlight_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {highlight}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(10)
            
            y_position += 1
    
    def _add_financials_slide(self, content: Dict):
        """Add financial overview slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Financial Overview')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Metrics grid
        metrics = content.get('metrics', [])
        x_positions = [1, 5.5]
        y_positions = [2, 4]
        
        for i, metric in enumerate(metrics[:4]):
            col = i % 2
            row = i // 2
            
            # Metric box with background
            metric_box = slide.shapes.add_shape(
                1,
                Inches(x_positions[col]), Inches(y_positions[row]),
                Inches(3.5), Inches(1.5)
            )
            metric_box.fill.solid()
            metric_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            metric_box.fill.fore_color.brightness = -0.15
            metric_box.line.fill.background()
            
            # Add text
            text_box = slide.shapes.add_textbox(
                Inches(x_positions[col] + 0.2), Inches(y_positions[row] + 0.2),
                Inches(3.1), Inches(1.1)
            )
            tf = text_box.text_frame
            
            # Label
            p1 = tf.paragraphs[0]
            p1.text = metric.get('label', 'Metric')
            p1.font.size = Pt(14)
            p1.font.color.rgb = RGBColor(200, 200, 200)
            
            # Value
            p2 = tf.add_paragraph()
            p2.text = metric.get('value', 'N/A')
            p2.font.size = Pt(28)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_breakdown_slide(self, content: Dict, slide_type: str):
        """Add assets or liabilities breakdown slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
            )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', f'{slide_type.title()} Breakdown')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Overview section
        overview_box = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(4), Inches(3)
            )
        tf = overview_box.text_frame
        
        p1 = tf.paragraphs[0]
        p1.text = "Overview"
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        
        items = [
            f"Total: {content.get('total', 'N/A')}",
            f"Current: {content.get('current', 'N/A')}",
            f"Non-Current: {content.get('non_current', 'N/A') or content.get('long_term', 'N/A')}"
        ]
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.space_before = Pt(10)
        
        # Breakdown section
        breakdown = content.get('breakdown', [])
        if breakdown:
            breakdown_box = slide.shapes.add_textbox(
                Inches(5.5), Inches(2),
                Inches(4), Inches(4)
            )
            tf2 = breakdown_box.text_frame
            
            p1 = tf2.paragraphs[0]
            p1.text = "Breakdown"
            p1.font.size = Pt(24)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(255, 255, 255)
            
            for item in breakdown[:5]:
                p = tf2.add_paragraph()
                p.text = f"{item.get('category', 'Item')}: {item.get('percentage', 0)}%"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_before = Pt(8)
    
    def _add_ratios_slide(self, content: Dict):
        """Add financial ratios slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
            )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
            )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Financial Ratios')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Ratios
        ratios = content.get('ratios', [])
        y_position = 1.8
        
        for ratio in ratios[:4]:
            # Ratio box
            ratio_box = slide.shapes.add_shape(
                1,
                Inches(1), Inches(y_position),
                Inches(8), Inches(1.2)
            )
            ratio_box.fill.solid()
            ratio_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            ratio_box.fill.fore_color.brightness = -0.15
            ratio_box.line.fill.background()
            
            # Add text
            text_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(y_position + 0.1),
                Inches(7.6), Inches(1)
            )
            tf = text_box.text_frame
            
            # Ratio name and value
            p1 = tf.paragraphs[0]
            p1.text = f"{ratio.get('name', 'Ratio')}: {ratio.get('value', 'N/A')}"
            p1.font.size = Pt(20)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(255, 255, 255)
            
            # Interpretation
            p2 = tf.add_paragraph()
            p2.text = ratio.get('interpretation', '')
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(200, 200, 200)
            
            y_position += 1.4
    
    def _add_insights_slide(self, content: Dict):
        """Add trends/insights or conclusion slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
            )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
            )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Insights')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        y_position = 1.8
        
        # Key Insights or Takeaways
        insights = content.get('insights', []) or content.get('key_takeaways', [])
        if insights:
            section_box = slide.shapes.add_textbox(
                Inches(1), Inches(y_position),
                Inches(8), Inches(2.5)
            )
            tf = section_box.text_frame
            
            p1 = tf.paragraphs[0]
            p1.text = "Key Insights" if 'insights' in content else "Key Takeaways"
            p1.font.size = Pt(24)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(255, 255, 255)
            
            for insight in insights[:4]:
                p = tf.add_paragraph()
                p.text = f"✓ {insight}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_before = Pt(10)
            
            y_position += 3
        
        # Next Steps or Recommendations
        next_steps = content.get('next_steps', []) or content.get('recommendations', [])
        if next_steps:
            section_box = slide.shapes.add_textbox(
                Inches(1), Inches(y_position),
                Inches(8), Inches(2)
                )
            tf = section_box.text_frame
            
            p1 = tf.paragraphs[0]
            p1.text = "Next Steps" if 'next_steps' in content else "Recommendations"
            p1.font.size = Pt(24)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(255, 255, 255)
            
            for i, step in enumerate(next_steps[:3], 1):
                p = tf.add_paragraph()
                p.text = f"{i}. {step}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_before = Pt(10)
    
    def _add_company_slide(self, content: Dict):
        """Add company profile slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
            )
        title_frame = title_box.text_frame
        title_frame.text = content.get('company_name', 'Company Profile')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Company info
        info_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Inches(8), Inches(5)
            )
        tf = info_box.text_frame
            
        # Industry and Founded
        if content.get('industry') or content.get('founded'):
            p = tf.paragraphs[0]
            info_text = []
            if content.get('industry'):
                info_text.append(f"Industry: {content['industry']}")
            if content.get('founded'):
                info_text.append(f"Founded: {content['founded']}")
            p.text = " | ".join(info_text)
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.space_after = Pt(20)
        
        # Mission
        if content.get('mission'):
            p = tf.add_paragraph()
            p.text = "Mission"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(20)
            
            p2 = tf.add_paragraph()
            p2.text = content['mission']
            p2.font.size = Pt(16)
            p2.font.color.rgb = RGBColor(200, 200, 200)
            p2.space_after = Pt(20)
        
        # Key Facts
        if content.get('key_facts'):
                p = tf.add_paragraph()
            p.text = "Key Facts"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(20)
            
            for fact in content['key_facts'][:4]:
                p = tf.add_paragraph()
                p.text = f"• {fact}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_before = Pt(5)
    
    def _add_products_services_slide(self, content: Dict):
        """Add products & services slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
            )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Products & Services')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Content box
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Inches(8), Inches(5)
            )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Products
        if content.get('products'):
            p = tf.paragraphs[0]
            p.text = "Products & Services:"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(12)
            
            for product in content['products']:
                p = tf.add_paragraph()
                p.text = f"• {product}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_after = Pt(8)
        
        # Categories
        if content.get('categories'):
            p = tf.add_paragraph()
            p.text = "Categories:"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(16)
            p.space_after = Pt(12)
            
            for category in content['categories']:
                p = tf.add_paragraph()
                p.text = f"• {category}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_after = Pt(6)
        
        # Certifications
        if content.get('certifications'):
            p = tf.add_paragraph()
            p.text = "Certifications:"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(12)
            p.space_after = Pt(8)
            
            certs_text = " | ".join(content['certifications'][:4])
            p = tf.add_paragraph()
            p.text = certs_text
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200, 200, 200)
    
    def _add_markets_locations_slide(self, content: Dict):
        """Add markets & locations slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Markets & Locations')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Content box
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Inches(8), Inches(5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Markets
        if content.get('markets'):
            p = tf.paragraphs[0]
            p.text = "Markets Served:"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(12)
            
            for market in content['markets']:
                p = tf.add_paragraph()
                p.text = f"• {market}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_after = Pt(8)
        
        # Locations
        if content.get('locations'):
            p = tf.add_paragraph()
            p.text = "Locations:"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(16)
            p.space_after = Pt(12)
            
            for location in content['locations']:
                p = tf.add_paragraph()
                p.text = f"• {location}"
            p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_after = Pt(6)
        
        # Manufacturing
        if content.get('manufacturing'):
            p = tf.add_paragraph()
            p.text = "Manufacturing Capabilities:"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(12)
            p.space_after = Pt(8)
            
            p = tf.add_paragraph()
            p.text = content['manufacturing']
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200, 200, 200)
    
    def _add_leadership_slide(self, content: Dict):
        """Add leadership & team slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Leadership & Team')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Content box
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Inches(8), Inches(5)
            )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # CEO Message
        if content.get('ceo_message_summary'):
            p = tf.paragraphs[0]
            p.text = "CEO's Message:"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(12)
            
                p = tf.add_paragraph()
            p.text = content['ceo_message_summary']
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.space_after = Pt(20)
        
        # Leadership Team
        if content.get('leadership'):
            p = tf.add_paragraph()
            p.text = "Leadership Team:"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(12)
            
            for leader in content['leadership']:
                p = tf.add_paragraph()
                p.text = f"• {leader}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_after = Pt(8)
    
    def _add_major_projects_slide(self, content: Dict):
        """Add major projects & clients slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
            )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Major Projects & Clients')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Content box
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Inches(8), Inches(5)
            )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Major Projects
        if content.get('projects'):
            p = tf.paragraphs[0]
            p.text = "Major Projects:"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(12)
            
            for project in content['projects']:
                p = tf.add_paragraph()
                p.text = f"• {project}"
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_after = Pt(8)
        
        # Clients
        if content.get('clients'):
            p = tf.add_paragraph()
            p.text = "Notable Clients:"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(16)
            p.space_after = Pt(12)
            
            clients_text = " | ".join(content['clients'][:8])
                p = tf.add_paragraph()
            p.text = clients_text
                p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200, 200, 200)
    
    def _add_vision_mission_slide(self, content: Dict):
        """Add vision & mission slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # Background
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0),
            Inches(Config.SLIDE_WIDTH), Inches(Config.SLIDE_HEIGHT)
            )
        background.fill.solid()
        background.fill.fore_color.rgb = self.theme['primary']
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = content.get('title', 'Vision, Mission & Values')
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Content box
        content_box = slide.shapes.add_textbox(
            Inches(1), Inches(1.8),
            Inches(8), Inches(5)
            )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        # Vision
        if content.get('vision'):
            p = tf.paragraphs[0]
            p.text = "Vision:"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(8)
            
            p = tf.add_paragraph()
            p.text = content['vision']
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.space_after = Pt(16)
        
        # Mission
        if content.get('mission'):
            p = tf.add_paragraph()
            p.text = "Mission:"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(8)
            
            p = tf.add_paragraph()
            p.text = content['mission']
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(200, 200, 200)
            p.space_after = Pt(16)
        
        # Values
        if content.get('values'):
            p = tf.add_paragraph()
            p.text = "Core Values:"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_after = Pt(8)
            
            for value in content['values']:
                p = tf.add_paragraph()
                p.text = f"• {value}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(200, 200, 200)
                p.space_after = Pt(6)
        
        # USPs
        if content.get('usps'):
            p = tf.add_paragraph()
            p.text = "Unique Selling Points:"
            p.font.size = Pt(18)
            p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(12)
            p.space_after = Pt(8)

            usps_text = " | ".join(content['usps'][:4])
            p = tf.add_paragraph()
            p.text = usps_text
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200, 200, 200)